#!/usr/bin/env python3
"""股权架构图生成工具 v5 - 支持多层架构/递归子公司/分界线/交叉持股/自定义间距"""
import io, json, os
from http.server import HTTPServer, BaseHTTPRequestHandler
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from lxml import etree

BLACK = RGBColor(0,0,0)
FONT_SZ = Pt(7)
BOX_BORDER = int(Pt(1).pt * 12700)
LINE_W = int(Pt(0.25).pt * 12700)
NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

BW_MIN, BH_MIN = 1.6, 1.0
BW_PAD, BH_PAD = 0.65, 0.45
LINE_H = 0.40
PW, PH = 1.5, 0.5
PCT_OFF = 0.12
SW, SH = 40.9, 19.1
NL = "\n"

_CJK = [('\u4e00','\u9fff'),('\u3000','\u303f'),('\uff00','\uffef'),('\u2e80','\u2eff')]
def _is_cjk(c): return any(lo<=c<=hi for lo,hi in _CJK)

def box_width(name):
    lines = [l for l in name.split(NL) if l.strip()]
    if not lines: return BW_MIN
    def lw(l): return sum(0.247 if _is_cjk(c) else 0.135 for c in l)
    return max(BW_MIN, max(lw(l) for l in lines) + BW_PAD)

def box_height(name):
    n = max(1, len([l for l in name.split(NL) if l]))
    return max(BH_MIN, n * LINE_H + BH_PAD)

def calc_positions(widths, gap, center=None):
    if center is None: center = SW / 2
    n = len(widths)
    if n == 0: return []
    g = gap
    total = sum(widths) + (n-1)*g
    usable = SW - 2.4
    if total > usable and center == SW/2:
        g = max(gap*0.3, (usable - sum(widths)) / max(n-1,1))
        total = sum(widths) + (n-1)*g
    x = center - total/2
    cxs = []
    for w in widths:
        cxs.append(x + w/2); x += w + g
    return cxs

def count_leaves(node):
    ch = node.get('children', [])
    return sum(count_leaves(c) for c in ch) if ch else 1

def tree_depth(nodes, cur=0):
    if not nodes: return cur
    return max(tree_depth(n.get('children',[]), cur+1) for n in nodes)

def get_level_max_h(nodes, level, cur=0):
    mh = BH_MIN
    for node in nodes:
        if cur == level: mh = max(mh, box_height(node['name']))
        if node.get('children') and cur < level:
            mh = max(mh, get_level_max_h(node['children'], level, cur+1))
    return mh

def get_level_max_w(nodes, level, cur=0):
    mw = BW_MIN
    for node in nodes:
        if cur == level: mw = max(mw, box_width(node['name']))
        if node.get('children') and cur < level:
            mw = max(mw, get_level_max_w(node['children'], level, cur+1))
    return mw

def normalize_level_dims(nodes, level_w, level_h, cur=0):
    for node in nodes:
        node['_w'] = level_w[cur]
        node['_h'] = level_h[cur]
        if node.get('children'):
            normalize_level_dims(node['children'], level_w, level_h, cur+1)

def assign_tree_compact(nodes, center, gap):
    ws = [n['_w'] for n in nodes]
    cxs = calc_positions(ws, gap, center)
    for i, node in enumerate(nodes):
        node['_cx'] = cxs[i]
        if node.get('children'):
            assign_tree_compact(node['children'], cxs[i], gap)

# ── PPTX helpers ──────────────────────────────────────────────────────────────
def _rm_shadow(shape):
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is not None and spPr.find(qn('a:effectLst')) is None:
        etree.SubElement(spPr, qn('a:effectLst'))

def _set_border(tb, w, dashed=False):
    sp = tb._element; spPr = sp.find(qn('p:spPr'))
    for e in spPr.findall(qn('a:ln')): spPr.remove(e)
    dash = '<a:prstDash xmlns:a="{}" val="dash"/>'.format(NS) if dashed else ''
    spPr.append(etree.fromstring(
        '<a:ln xmlns:a="{}" w="{}"><a:solidFill><a:srgbClr val="000000"/>'
        '</a:solidFill>{}</a:ln>'.format(NS, w, dash)))

def _set_text(tf, lines, align=PP_ALIGN.CENTER):
    bp = tf._txBody.find(qn('a:bodyPr'))
    bp.set('anchor','ctr'); bp.set('lIns','0'); bp.set('rIns','0')
    tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align
        if line.strip():
            run = p.add_run()
            run.text = line; run.font.size=FONT_SZ; run.font.color.rgb=BLACK
            run.font.name = 'Times New Roman'
            rPr = run._r.find(qn('a:rPr'))
            if rPr is None: rPr = etree.SubElement(run._r, qn('a:rPr'))
            for ea in rPr.findall(qn('a:ea')): rPr.remove(ea)
            ea = etree.SubElement(rPr, qn('a:ea')); ea.set('typeface','微软雅黑')

def add_box(slide, x, y, w, h, lines, dashed=False, filled=False):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    _set_border(tb, LINE_W if dashed else BOX_BORDER, dashed)
    if filled:
        spPr = tb._element.find(qn('p:spPr'))
        spPr.append(etree.fromstring(
            '<a:solidFill xmlns:a="{}"><a:srgbClr val="D9D9D9"/></a:solidFill>'.format(NS)))
    _set_text(tb.text_frame, lines if lines else [''])

def add_label(slide, x, y, w, h, text):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    sp = tb._element; spPr = sp.find(qn('p:spPr'))
    for e in spPr.findall(qn('a:ln')): spPr.remove(e)
    _set_text(tb.text_frame, [text], align=PP_ALIGN.LEFT)

def add_xl_label(slide, x, y, text):
    """Narrow centered label for cross-link percentages"""
    W, H = 0.8, 0.5
    tb = slide.shapes.add_textbox(Cm(x - W/2), Cm(y), Cm(W), Cm(H))
    sp = tb._element; spPr = sp.find(qn('p:spPr'))
    for e in spPr.findall(qn('a:ln')): spPr.remove(e)
    bp = tb.text_frame._txBody.find(qn('a:bodyPr'))
    bp.set('anchor','ctr'); bp.set('lIns','18000'); bp.set('rIns','18000')
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size = FONT_SZ; run.font.color.rgb = BLACK
    run.font.name = 'Times New Roman'
    rPr = run._r.find(qn('a:rPr'))
    if rPr is None: rPr = etree.SubElement(run._r, qn('a:rPr'))
    for ea in rPr.findall(qn('a:ea')): rPr.remove(ea)
    ea = etree.SubElement(rPr, qn('a:ea')); ea.set('typeface','微软雅黑')

def add_line(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1),Cm(y1),Cm(x2),Cm(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(0.25); _rm_shadow(c)

def add_arrow(slide, x, y1, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x),Cm(y1),Cm(x),Cm(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(0.25); _rm_shadow(c)
    ln = c._element.find(qn('p:spPr')).find(qn('a:ln'))
    if ln is not None:
        ln.append(etree.fromstring('<a:tailEnd xmlns:a="{}" type="triangle" w="sm" len="sm"/>'.format(NS)))

def add_harrow(slide, x1, y, x2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1),Cm(y),Cm(x2),Cm(y))
    c.line.color.rgb = BLACK; c.line.width = Pt(0.25); _rm_shadow(c)
    ln = c._element.find(qn('p:spPr')).find(qn('a:ln'))
    if ln is not None:
        ln.append(etree.fromstring('<a:tailEnd xmlns:a="{}" type="triangle" w="sm" len="sm"/>'.format(NS)))

def add_hline_dashed(slide, x1, y, x2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1),Cm(y),Cm(x2),Cm(y))
    c.line.color.rgb = BLACK; c.line.width = Pt(0.5); _rm_shadow(c)
    ln = c._element.find(qn('p:spPr')).find(qn('a:ln'))
    if ln is not None:
        ln.append(etree.fromstring('<a:prstDash xmlns:a="{}" val="dash"/>'.format(NS)))

# ── Main generator ────────────────────────────────────────────────────────────
def generate_pptx(data):
    prs = Presentation()
    prs.slide_width = Cm(SW); prs.slide_height = Cm(SH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shs   = data['shareholders']
    subs  = data['subsidiaries']
    comp  = data.get('companyName', '本公司')
    mid   = data.get('middleTier', [])
    div   = data.get('divider', None)
    xlink = data.get('crossLinks', [])
    pcs   = data.get('parentControllers', [])

    # Spacing (frontend overrides)
    sp    = data.get('spacing', {})
    SH_GAP  = float(sp.get('shGap',  0.8))   # shareholder horizontal gap
    PC_GAP  = float(sp.get('pcGap',  0.8))   # parent-controller horizontal gap
    SUB_GAP = float(sp.get('subGap', 0.8))   # subsidiary horizontal gap
    GAP_V   = float(sp.get('gapV',   0.75))  # vertical gap between layers

    PC_H = 0.85; PC_STUB = 0.4; PC_BUS = 0.5

    # Shareholder sizing
    sh_w = [box_width(s['name']) for s in shs]
    sh_h = [box_height(s['name']) for s in shs]
    msh = max(sh_h); msw = max(sh_w)
    sh_cx = calc_positions([msw]*len(shs), SH_GAP)

    comp_w = box_width(comp); comp_h = box_height(comp)

    # ── Vertical centering ────────────────────────────────────────────────────
    depth_pre = tree_depth(subs)
    lh_pre = [get_level_max_h(subs, lv) for lv in range(depth_pre)] if depth_pre else []
    pc_above = (PC_H + PC_STUB + PC_BUS) if pcs else 0
    h_content = pc_above + msh + GAP_V
    if mid:
        h_content += GAP_V + max(box_height(m['name']) for m in mid) + GAP_V
    h_content += comp_h + GAP_V
    if lh_pre:
        h_content += GAP_V + sum(lh + GAP_V*2 for lh in lh_pre)
    R1Y = max(pc_above + 0.8, (SH - h_content) / 2 + pc_above)

    # ── Parent controllers (T-connector) ─────────────────────────────────────
    BUS_Y = R1Y - PC_BUS
    BOX_BOT = BUS_Y - PC_STUB
    for pc in pcs:
        pc_name = pc.get('name','').strip()
        ctrl_shs = pc.get('shareholders', [])
        if not pc_name or not ctrl_shs: continue
        idxs = [int(s['idx']) for s in ctrl_shs if int(s['idx']) < len(sh_cx)]
        if not idxs: continue
        pc_w = box_width(pc_name)
        cx = (min(sh_cx[i] for i in idxs) + max(sh_cx[i] for i in idxs)) / 2
        add_box(slide, cx-pc_w/2, BOX_BOT-PC_H, pc_w, PC_H, [pc_name])
        add_line(slide, cx, BOX_BOT, cx, BUS_Y)
        xs = [sh_cx[i] for i in idxs]
        if len(xs) > 1:
            add_line(slide, min(xs), BUS_Y, max(xs), BUS_Y)
        for j, s in enumerate(ctrl_shs):
            idx = int(s['idx'])
            if idx >= len(sh_cx): continue
            add_arrow(slide, sh_cx[idx], BUS_Y, R1Y)
            if s.get('pct','').strip():
                add_label(slide, sh_cx[idx]+PCT_OFF, BUS_Y+(PC_BUS-PH)/2, PW, PH, s['pct'])

    # ── Shareholders ──────────────────────────────────────────────────────────
    for i, sh in enumerate(shs):
        add_box(slide, sh_cx[i]-msw/2, R1Y, msw, msh, sh['name'].split(NL))

    ctrl = [i for i,s in enumerate(shs) if s.get('isControl')]
    if ctrl:
        lx = sh_cx[ctrl[0]]  - msw/2 - 0.3
        rx = sh_cx[ctrl[-1]] + msw/2 + 0.3
        add_box(slide, lx, R1Y-0.3, rx-lx, msh+0.6, [], dashed=True)

    CCX = SW / 2
    cur_y = R1Y + msh

    # ── Middle tier ───────────────────────────────────────────────────────────
    if mid:
        mid_w = [box_width(m['name']) for m in mid]
        mid_h = [box_height(m['name']) for m in mid]
        mid_cx = calc_positions(mid_w, SH_GAP)
        mmid = max(mid_h); mmw = max(mid_w)
        CONN_Y  = cur_y + GAP_V * 0.5    # T-junction: parents → mid-tier
        MID_Y   = CONN_Y + GAP_V * 0.5   # mid-tier box top
        MID_BOT = MID_Y + mmid            # mid-tier box bottom
        CY      = MID_BOT + GAP_V         # 本公司 Y (also direct-sh merge point)
        mid_flat = set(pi for m in mid for pi in m.get('parentShareholders',[]))
        # Draw mid-tier boxes and parent-shareholder T-connections
        for i, m in enumerate(mid):
            add_box(slide, mid_cx[i]-mmw/2, MID_Y, mmw, mmid, m['name'].split(NL))
            # Mid-tier bottom → CY merge bus
            add_line(slide, mid_cx[i], MID_BOT, mid_cx[i], CY)
            add_label(slide, mid_cx[i]+PCT_OFF, MID_BOT+(CY-MID_BOT)/2-PH/2, PW, PH, m.get('mainPct',''))
            # Parent shareholders T-junction at CONN_Y
            pis = [pi for pi in m.get('parentShareholders', []) if 0 <= pi < len(sh_cx)]
            if pis:
                xs = [sh_cx[pi] for pi in pis]
                for pi in pis:
                    add_line(slide, sh_cx[pi], cur_y, sh_cx[pi], CONN_Y)
                    add_label(slide, sh_cx[pi]+PCT_OFF, cur_y+(CONN_Y-cur_y)/2-PH/2, PW, PH, shs[pi].get('pct',''))
                # Horizontal bus connecting all parents
                if len(xs) > 1:
                    add_line(slide, min(xs), CONN_Y, max(xs), CONN_Y)
                bus_x = (min(xs) + max(xs)) / 2
                # Bus → mid-tier top arrow
                if abs(bus_x - mid_cx[i]) > 0.05:
                    add_line(slide, bus_x, CONN_Y, mid_cx[i], CONN_Y)
                add_arrow(slide, mid_cx[i], CONN_Y, MID_Y)
        # Non-participating shareholders: draw direct line all the way to CY
        for i, sh in enumerate(shs):
            if i not in mid_flat:
                add_line(slide, sh_cx[i], cur_y, sh_cx[i], CY)
                add_label(slide, sh_cx[i]+PCT_OFF, cur_y+(CY-cur_y)/2-PH/2, PW, PH, sh['pct'])
        # Horizontal merge bus at CY joining all paths
        all_x = [sh_cx[i] for i in range(len(shs)) if i not in mid_flat] + list(mid_cx)
        if len(all_x) > 1: add_line(slide, min(all_x), CY, max(all_x), CY)
        add_arrow(slide, CCX, CY, CY + GAP_V)
        CY = CY + GAP_V
    else:
        MERGE_Y = cur_y + GAP_V
        CY = MERGE_Y + GAP_V
        for i, sh in enumerate(shs):
            add_line(slide, sh_cx[i], cur_y, sh_cx[i], MERGE_Y)
            add_label(slide, sh_cx[i]+PCT_OFF, cur_y+(MERGE_Y-cur_y)/2-PH/2, PW, PH, sh['pct'])
        if len(sh_cx) > 1: add_line(slide, sh_cx[0], MERGE_Y, sh_cx[-1], MERGE_Y)
        add_arrow(slide, CCX, MERGE_Y, CY)

    # ── Company box ───────────────────────────────────────────────────────────
    filled = data.get('companyFilled', False)
    add_box(slide, CCX-comp_w/2, CY, comp_w, comp_h, [comp], filled=filled)

    if not subs:
        buf = io.BytesIO(); prs.save(buf); return buf.getvalue()

    # ── Subsidiaries ──────────────────────────────────────────────────────────
    depth = tree_depth(subs)
    level_h = [get_level_max_h(subs, lv) for lv in range(depth)]
    level_w = [get_level_max_w(subs, lv) for lv in range(depth)]
    normalize_level_dims(subs, level_w, level_h)
    assign_tree_compact(subs, SW/2, SUB_GAP)

    SM = CY + comp_h + GAP_V
    level_y = []
    y = SM + GAP_V
    level_y.append(y)
    for lv in range(1, depth):
        y += level_h[lv-1] + GAP_V*2
        level_y.append(y)

    div_y = None
    if div:
        after_lv = div.get('afterLevel', 0)
        if after_lv < len(level_y):
            div_y = level_y[after_lv] + level_h[after_lv] + GAP_V * 0.6

    add_line(slide, CCX, CY+comp_h, CCX, SM)
    if len(subs) > 1:
        add_line(slide, min(n['_cx'] for n in subs), SM, max(n['_cx'] for n in subs), SM)

    box_pos = {}

    def draw_tree(nodes, parent_sm, level):
        mh = level_h[level]; mw = level_w[level]; SY = level_y[level]
        for node in nodes:
            cx = node['_cx']
            add_arrow(slide, cx, parent_sm, SY)
            add_label(slide, cx+PCT_OFF, parent_sm+(SY-parent_sm)/2-PH/2, PW, PH, node['pct'])
            add_box(slide, cx-mw/2, SY, mw, mh, node['name'].split(NL))
            box_pos[node['name'].replace(NL,'')] = {'cx':cx,'w':mw,'y':SY,'h':mh}
            if node.get('children'):
                ch = node['children']
                ch_sm = SY + mh + GAP_V
                add_line(slide, cx, SY+mh, cx, ch_sm)
                if len(ch) > 1:
                    add_line(slide, min(c['_cx'] for c in ch), ch_sm, max(c['_cx'] for c in ch), ch_sm)
                draw_tree(ch, ch_sm, level+1)

    draw_tree(subs, SM, 0)

    if div_y:
        add_hline_dashed(slide, 1.0, div_y, SW-1.0)
        add_label(slide, 1.2, div_y-0.35, 2.0, 0.4, div.get('leftLabel','境外'))
        add_label(slide, SW-3.2, div_y+0.05, 2.0, 0.4, div.get('rightLabel','境内'))

    # Cross-links: horizontal arrow
    for xl in xlink:
        fn = xl.get('from','').replace(NL,'')
        tn = xl.get('to','').replace(NL,'')
        if fn in box_pos and tn in box_pos:
            fp = box_pos[fn]; tp = box_pos[tn]
            mid_y = (fp['y'] + fp['h']/2 + tp['y'] + tp['h']/2) / 2
            if fp['cx'] < tp['cx']:
                x1 = fp['cx'] + fp['w']/2; x2 = tp['cx'] - tp['w']/2
            else:
                x1 = fp['cx'] - fp['w']/2; x2 = tp['cx'] + tp['w']/2
            add_harrow(slide, x1, mid_y, x2)
            if xl.get('pct','').strip():
                add_xl_label(slide, (x1+x2)/2, mid_y-PH-0.05, xl['pct'])

    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


# ── HTTP server ───────────────────────────────────────────────────────────────
HTML = """<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8"/><title>股权架构图生成工具</title>
<style>
@import url('https://fonts.googleapis.com/css2?family=Noto+Serif+SC:wght@400;700&family=Noto+Sans+SC:wght@400;500;700&display=swap');
*{box-sizing:border-box;margin:0;padding:0;}
body{font-family:'Noto Sans SC',system-ui,-apple-system,sans-serif;background:#f0f2f8;}
header{background:#1a2640;color:#fff;padding:14px 28px;display:flex;align-items:center;gap:12px;position:sticky;top:0;z-index:100;}
header h1{font-size:20px;font-weight:600;letter-spacing:.5px;}
.layout{display:flex;gap:0;height:calc(100vh - 52px);overflow:hidden;position:relative;}
.panel{width:420px;min-width:260px;max-width:640px;background:#f7f8fc;border-right:1px solid #e2e6ef;padding:20px 16px;overflow-y:auto;flex-shrink:0;height:100%;}
.resize-handle{width:5px;background:#e2e6ef;cursor:col-resize;flex-shrink:0;transition:background .15s;}
.resize-handle:hover,.resize-handle.dragging{background:#a0b0cc;}
.preview-area{flex:1;padding:20px;display:flex;flex-direction:column;gap:12px;overflow-y:auto;height:100%;min-width:0;}
.card{background:#fff;border:1px solid #e2e6ef;border-radius:10px;padding:14px 16px;margin-bottom:12px;}
.card-title{font-size:15px;font-weight:600;color:#1a2640;margin-bottom:10px;display:flex;justify-content:space-between;align-items:center;}
.section-toggle{font-size:12px;color:#7a8ba6;cursor:pointer;font-weight:400;user-select:none;padding:2px 8px;border:1px solid #dde3ed;border-radius:12px;}
.section-toggle:hover{background:#f0f2f8;}
.extra-section{display:none;}
.hint{font-size:13px;color:#7a8ba6;margin-bottom:8px;}
input[type=text],input[type=number],textarea,select{
  width:100%;padding:6px 9px;border:1px solid #dde3ed;border-radius:6px;
  font-size:13px;font-family:inherit;background:#fff;color:#1a2640;}
input[type=checkbox]{width:16px;height:16px;accent-color:#1a2640;cursor:pointer;}
textarea{resize:vertical;min-height:36px;}
.row{display:flex;gap:8px;align-items:center;margin-bottom:6px;}
.row label{font-size:13px;color:#444;white-space:nowrap;}
.sh-row,.sub-row,.mid-row,.xl-row,.pc-sh-row{
  display:flex;gap:6px;align-items:flex-start;margin-bottom:6px;padding:8px;
  background:#f7f8fc;border-radius:7px;border:1px solid #e8ecf4;}
.sh-row .pi,.sub-row .pi,.mid-row .pi,.xl-row .pi{width:70px;flex-shrink:0;}
.rm{background:none;border:none;color:#c0c8d8;cursor:pointer;font-size:16px;padding:2px 4px;line-height:1;}
.rm:hover{color:#e05;}
.add-btn{font-size:13px;color:#1a2640;border:1px dashed #b0bcd4;background:#f7f8fc;
  border-radius:6px;padding:5px 12px;cursor:pointer;width:100%;margin-top:4px;}
.add-btn:hover{background:#e8edf7;}
.gen-wrap{text-align:center;margin:24px 0 8px;}
#genBtn{background:#1a2640;color:#fff;border:none;border-radius:26px;
  padding:14px 44px;font-size:16px;font-family:inherit;cursor:pointer;
  box-shadow:0 2px 10px rgba(26,38,64,.25);transition:.15s;}
#genBtn:hover{background:#253659;}
#genBtn:disabled{opacity:.6;cursor:not-allowed;}
#status{font-size:13px;margin-top:8px;min-height:18px;}
.ok{color:#2a7a3b;} .err{color:#c0392b;}
.preview-card{background:#fff;border:1px solid #e2e6ef;border-radius:10px;padding:14px 16px;}
.preview-title{font-size:14px;font-weight:600;color:#1a2640;margin-bottom:2px;}
.preview-sub{font-size:12px;color:#7a8ba6;margin-bottom:10px;}
#svgWrap{overflow:auto;border:1px solid #e8ecf4;border-radius:6px;min-height:120px;background:#fafbff;display:flex;align-items:flex-start;justify-content:center;}
/* Spacing row */
.sp-row{display:flex;align-items:center;gap:10px;margin-bottom:10px;}
.sp-label{font-size:13px;color:#333;flex:0 0 150px;}
.sp-row input[type=range]{flex:1;accent-color:#1a2640;height:4px;}
.sp-row input[type=number]{width:62px;flex-shrink:0;text-align:center;}
.sp-unit{font-size:12px;color:#aaa;flex-shrink:0;}
</style>
</head>
<body>
<header><h1>📊 股权架构图生成工具</h1></header>
<div class="layout">
<div class="panel" id="mainPanel">

<!-- 本公司 -->
<div class="card">
  <div class="card-title">本公司</div>
  <div class="row"><label>名称</label>
    <input type="text" id="companyName" value="本公司" oninput="safeRenderPreview()"/>
  </div>
  <div class="row"><label><input type="checkbox" id="companyFilled" onchange="safeRenderPreview()"> 灰色填充</label></div>
</div>

<!-- 股东 -->
<div class="card">
  <div class="card-title">股东层</div>
  <div id="shareholderList"></div>
  <button class="add-btn" onclick="addShareholder()">+ 添加股东</button>
</div>

<!-- 上层控股人 -->
<div class="card">
  <div class="card-title">上层控股人（可选）
    <span class="section-toggle" onclick="toggleSection('pcSection')">▶ 展开</span>
  </div>
  <div id="pcSection" class="extra-section">
    <p class="hint">控股人持有一个或多个股东，箭头从上层直指股东</p>
    <div id="pcList"></div>
    <button class="add-btn" onclick="addPC()">+ 添加上层控股人</button>
  </div>
</div>

<!-- 中间控股层 -->
<div class="card">
  <div class="card-title">中间控股层（可选）
    <span class="section-toggle" onclick="toggleSection('midSection')">▶ 展开</span>
  </div>
  <div id="midSection" class="extra-section">
    <p class="hint">股东通过中间公司控制本公司（如开曼控股层）</p>
    <div id="midList"></div>
    <button class="add-btn" onclick="addMid()">+ 添加中间层公司</button>
  </div>
</div>

<!-- 子公司 -->
<div class="card">
  <div class="card-title">子公司层</div>
  <div id="subsidiaryList"></div>
  <button class="add-btn" onclick="addSubsidiary(null)">+ 添加子公司</button>
</div>

<!-- 交叉持股 -->
<div class="card">
  <div class="card-title">交叉持股（可选）
    <span class="section-toggle" onclick="toggleSection('xlSection')">▶ 展开</span>
  </div>
  <div id="xlSection" class="extra-section">
    <div id="xlList"></div>
    <button class="add-btn" onclick="addXL()">+ 添加交叉持股</button>
  </div>
</div>

<!-- 间距调整 -->
<div class="card">
  <div class="card-title">间距调整
    <span class="section-toggle" onclick="toggleSection('spacingSection')">▶ 展开</span>
  </div>
  <div id="spacingSection" class="extra-section">
    <p class="hint">分别调整各层方框的水平/垂直间距（cm）</p>
    <div id="spacingList"></div>
  </div>
</div>

<div class="gen-wrap">
  <button id="genBtn" onclick="generate()">&#8595; 生成并下载 PPTX</button>
  <p id="status"></p>
</div>

</div><!-- /panel -->
<div class="resize-handle" id="resizeHandle"></div>
<div class="preview-area">
  <div class="preview-card">
    <div class="preview-title" style="display:flex;align-items:center;justify-content:space-between;margin-bottom:6px">
    <span>实时预览 <span style="font-size:12px;font-weight:400;color:#7a8ba6">布局与导出PPT完全一致</span></span>
    <span style="display:flex;gap:4px;align-items:center">
      <button onclick="previewZoom=Math.max(0.2,previewZoom===0?0.8:previewZoom-0.2);safeRenderPreview()" style="padding:2px 8px;border:1px solid #ccc;border-radius:4px;background:#fff;cursor:pointer;font-size:13px" title="缩小">－</button>
      <button onclick="previewZoom=0;safeRenderPreview()" style="padding:2px 8px;border:1px solid #ccc;border-radius:4px;background:#fff;cursor:pointer;font-size:12px" title="自适应">适应</button>
      <button onclick="previewZoom=(previewZoom===0?1.0:previewZoom)+0.2;safeRenderPreview()" style="padding:2px 8px;border:1px solid #ccc;border-radius:4px;background:#fff;cursor:pointer;font-size:13px" title="放大">＋</button>
    </span>
  </div>
    <div id="svgWrap"></div>
  </div>
</div>
</div><!-- /layout -->

<script>
// ── State ─────────────────────────────────────────────────────────────────────
let shareholders=[
  {name:"股东A", pct:"60%", isControl:true},
  {name:"股东B", pct:"40%", isControl:false},
];
let parentControllers=[];
let middleTier=[];
let subsidiaries=[
  {name:"子公司a", pct:"100%", children:[]},
  {name:"子公司b", pct:"100%", children:[]},
];
let crossLinks=[];

// Spacing state (cm)
let previewZoom = 0; // 0 = auto-fit, >0 = fixed scale
let spacing={
  shGap:  0.8,   // 股东层水平间距
  pcGap:  0.8,   // 上层控股人水平间距
  subGap: 0.8,   // 子公司水平间距
  gapV:   0.75,  // 各层垂直间距
};

// ── Helpers ───────────────────────────────────────────────────────────────────
function esc(s){return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');}
function autoH(ta){ta.style.height='auto';ta.style.height=(ta.scrollHeight)+'px';}
function toggleSection(id){
  const el=document.getElementById(id);
  // Find toggle span inside previousElementSibling (card-title), NOT replace its text
  const toggler=el.previousElementSibling&&el.previousElementSibling.querySelector('.section-toggle');
  if(el.style.display==='block'){
    el.style.display='none';
    if(toggler)toggler.textContent='▶ 展开';
  } else {
    el.style.display='block';
    if(toggler)toggler.textContent='▼ 收起';
  }
}
function openSection(id){
  const el=document.getElementById(id);
  el.style.display='block';
  const toggler=el.previousElementSibling&&el.previousElementSibling.querySelector('.section-toggle');
  if(toggler)toggler.textContent='▼ 收起';
}
function getAllNodeNames(nodes){
  let names=[];
  nodes.forEach(n=>{names.push(n.name.replace(/\\n/g,''));if(n.children)names.push(...getAllNodeNames(n.children));});
  return names;
}

// ── Render UI lists ───────────────────────────────────────────────────────────
function renderShareholders(){
  document.getElementById('shareholderList').innerHTML=shareholders.map((sh,i)=>`
    <div class="sh-row">
      <div style="flex:1;display:flex;flex-direction:column;gap:4px">
        <textarea rows="1" placeholder="股东名称（换行=\\n）" oninput="shareholders[${i}].name=this.value;autoH(this);renderPreview()" style="width:100%">${esc(sh.name)}</textarea>
        <div style="display:flex;gap:6px">
          <input class="pi" type="text" value="${esc(sh.pct)}" placeholder="持股%" oninput="shareholders[${i}].pct=this.value;renderPreview()"/>
          <label style="font-size:12px;display:flex;align-items:center;gap:4px;white-space:nowrap">
            <input type="checkbox" ${sh.isControl?'checked':''} onchange="shareholders[${i}].isControl=this.checked;renderPreview()"> 控股股东
          </label>
        </div>
      </div>
      <button class="rm" onclick="shareholders.splice(${i},1);renderShareholders();renderPreview()">&#x2715;</button>
    </div>`).join('');
}
function addShareholder(){shareholders.push({name:'',pct:'',isControl:false});renderShareholders();}

function renderParentControllers(){
  document.getElementById('pcList').innerHTML=parentControllers.map((pc,pi)=>`
    <div class="sh-row">
      <div style="flex:1;display:flex;flex-direction:column;gap:6px">
        <div style="display:flex;gap:6px;align-items:center">
          <textarea rows="1" style="flex:1" placeholder="控股人名称" oninput="parentControllers[${pi}].name=this.value;autoH(this);renderPreview()">${esc(pc.name)}</textarea>
          <button class="rm" onclick="parentControllers.splice(${pi},1);renderParentControllers();renderPreview()">&#x2715;</button>
        </div>
        <div style="font-size:12px;color:#666;margin-bottom:2px">控股的股东：</div>
        ${(pc.shareholders||[]).map((s,si)=>`
          <div class="pc-sh-row" style="display:flex;gap:6px;align-items:center">
            <select style="flex:1;padding:5px 6px;border:1px solid #ddd;border-radius:4px;font-size:12px" onchange="parentControllers[${pi}].shareholders[${si}].idx=+this.value;renderPreview()">
              ${shareholders.map((sh,i)=>`<option value="${i}" ${s.idx===i?'selected':''}>${esc(sh.name.replace(/\\n/g,' '))}</option>`).join('')}
            </select>
            <input type="text" value="${esc(s.pct)}" placeholder="持股%" style="width:70px" oninput="parentControllers[${pi}].shareholders[${si}].pct=this.value;renderPreview()"/>
            <button class="rm" onclick="parentControllers[${pi}].shareholders.splice(${si},1);renderParentControllers();renderPreview()">&#x2715;</button>
          </div>`).join('')}
        <button class="add-btn" onclick="parentControllers[${pi}].shareholders.push({idx:0,pct:''});renderParentControllers()">+ 添加控股的股东</button>
      </div>
    </div>`).join('');
}
function addPC(){parentControllers.push({name:'',shareholders:[{idx:0,pct:''}]});renderParentControllers();openSection('pcSection');}

function renderMid(){
  document.getElementById('midList').innerHTML=middleTier.map((m,i)=>`
    <div class="mid-row">
      <div style="flex:1;display:flex;flex-direction:column;gap:4px">
        <textarea rows="1" placeholder="公司名称" oninput="middleTier[${i}].name=this.value;autoH(this);renderPreview()" style="width:100%">${esc(m.name)}</textarea>
        <div style="display:flex;gap:6px">
          <input class="pi" type="text" value="${esc(m.mainPct||'')}" placeholder="持股%" oninput="middleTier[${i}].mainPct=this.value;renderPreview()"/>
          <select style="flex:1;padding:5px;border:1px solid #ddd;border-radius:4px;font-size:12px" onchange="middleTier[${i}].parentShareholders=[...this.selectedOptions].map(o=>+o.value);renderPreview()" multiple>
            ${shareholders.map((sh,si)=>`<option value="${si}" ${(m.parentShareholders||[]).includes(si)?'selected':''}>${esc(sh.name.replace(/\\n/g,' '))}</option>`).join('')}
          </select>
        </div>
      </div>
      <button class="rm" onclick="middleTier.splice(${i},1);renderMid();renderPreview()">&#x2715;</button>
    </div>`).join('');
}
function addMid(){middleTier.push({name:'',mainPct:'',parentShareholders:[]});renderMid();openSection('midSection');}

function renderSubsidiariesUI(nodes, parentPath, container){
  container.innerHTML='';
  nodes.forEach((sub,i)=>{
    const path=parentPath===null?i:`${parentPath}-${i}`;
    const div=document.createElement('div');
    div.className='sub-row'; div.style.flexDirection='column'; div.style.gap='6px';
    div.innerHTML=`
      <div style="display:flex;gap:6px;align-items:flex-start">
        <textarea rows="1" style="flex:1" placeholder="公司名称（换行=\\n）" oninput="getSubNode('${path}').name=this.value;autoH(this);renderPreview()">${esc(sub.name)}</textarea>
        <input class="pi" type="text" value="${esc(sub.pct)}" placeholder="持股%" oninput="getSubNode('${path}').pct=this.value;renderPreview()"/>
        <button class="rm" onclick="removeSubNode('${path}');renderSubsidiariesUI(subsidiaries,null,document.getElementById('subsidiaryList'));renderPreview()">&#x2715;</button>
      </div>
      <div style="display:flex;gap:8px">
        <button class="add-btn" style="flex:1" onclick="getSubNode('${path}').children.push({name:'',pct:'100%',children:[]});renderSubsidiariesUI(subsidiaries,null,document.getElementById('subsidiaryList'));renderPreview()">+ 添加子级</button>
      </div>`;
    container.appendChild(div);
    if(sub.children&&sub.children.length){
      const child=document.createElement('div');
      child.style.cssText='margin-left:18px;border-left:2px solid #e2e6ef;padding-left:10px;';
      renderSubsidiariesUI(sub.children, path, child);
      container.appendChild(child);
    }
  });
}
function getSubNode(path){
  const parts=String(path).split('-').map(Number);
  let node={children:subsidiaries};
  for(const p of parts) node=node.children[p];
  return node;
}
function removeSubNode(path){
  const parts=String(path).split('-').map(Number);
  const last=parts.pop();
  let parent={children:subsidiaries};
  for(const p of parts) parent=parent.children[p];
  parent.children.splice(last,1);
}
function addSubsidiary(){subsidiaries.push({name:'',pct:'100%',children:[]});renderSubsidiariesUI(subsidiaries,null,document.getElementById('subsidiaryList'));renderPreview();}

function renderXL(){
  const names=getAllNodeNames(subsidiaries);
  const opts=names.map(n=>`<option value="${esc(n)}">${esc(n)}</option>`).join('');
  document.getElementById('xlList').innerHTML=crossLinks.map((xl,i)=>`
    <div class="xl-row">
      <span style="font-size:12px;color:#555">从</span>
      <select style="flex:1;padding:5px 6px;border:1px solid #ddd;border-radius:4px;font-size:12px" onchange="crossLinks[${i}].from=this.value;renderPreview()">
        <option value="">-- 选择 --</option>${opts.replace(`value="${esc(xl.from)}"`,`value="${esc(xl.from)}" selected`)}
      </select>
      <span style="font-size:12px;color:#555">→</span>
      <select style="flex:1;padding:5px 6px;border:1px solid #ddd;border-radius:4px;font-size:12px" onchange="crossLinks[${i}].to=this.value;renderPreview()">
        <option value="">-- 选择 --</option>${opts.replace(`value="${esc(xl.to)}"`,`value="${esc(xl.to)}" selected`)}
      </select>
      <input class="pi" type="text" value="${esc(xl.pct)}" placeholder="持股%" oninput="crossLinks[${i}].pct=this.value;renderPreview()"/>
      <button class="rm" onclick="crossLinks.splice(${i},1);renderXL();renderPreview()">&#x2715;</button>
    </div>`).join('');
}
function addXL(){crossLinks.push({from:'',to:'',pct:''});renderXL();openSection('xlSection');}

function renderSpacing(){
  const items=[
    {key:'shGap',  label:'股东层水平间距',      min:0.1,max:4,step:0.1},
    {key:'pcGap',  label:'上层控股人水平间距',   min:0.1,max:4,step:0.1},
    {key:'subGap', label:'子公司层水平间距',     min:0.1,max:4,step:0.1},
    {key:'gapV',   label:'各层垂直间距',         min:0.1,max:3,step:0.1},
  ];
  document.getElementById('spacingList').innerHTML=items.map(it=>`
    <div class="sp-row">
      <span class="sp-label">${it.label}</span>
      <input type="range" min="${it.min}" max="${it.max}" step="${it.step}" value="${spacing[it.key]}"
        oninput="spacing['${it.key}']=+this.value;this.nextElementSibling.value=(+this.value).toFixed(1);renderPreview()"/>
      <input type="number" min="${it.min}" max="${it.max}" step="${it.step}" value="${spacing[it.key]}"
        oninput="spacing['${it.key}']=+this.value;this.previousElementSibling.value=this.value;renderPreview()"/>
      <span class="sp-unit">cm</span>
    </div>`).join('');
}

// ── SVG preview ───────────────────────────────────────────────────────────────
const BW_MIN=1.6,BH_MIN=1.0,BW_PAD=0.65,BH_PAD=0.45,LINE_H=0.40;
const SW=40.9,SH=19.1,POFF=0.12,PH=0.5,PW=1.5;

const _CJK=/[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef\u2e80-\u2eff]/;
function boxW(name){
  const ls=name.split(/\\n/).filter(l=>l.trim());
  if(!ls.length)return BW_MIN;
  const lw=l=>[...l].reduce((a,c)=>a+(_CJK.test(c)?0.247:0.135),0);
  return Math.max(BW_MIN,Math.max(...ls.map(lw))+BW_PAD);
}
function boxH(name){return Math.max(BH_MIN,Math.max(1,name.split(/\\n/).filter(l=>l).length)*LINE_H+BH_PAD);}

function calcPos(widths,gap,center){
  if(center===undefined)center=SW/2;
  const n=widths.length; if(!n)return [];
  const usable=SW-2.4;
  let g=gap, total=widths.reduce((a,w)=>a+w,0)+(n-1)*g;
  if(total>usable&&center===SW/2) g=Math.max(gap*0.3,(usable-widths.reduce((a,w)=>a+w,0))/Math.max(n-1,1));
  let x=center-(widths.reduce((a,w)=>a+w,0)+(n-1)*g)/2;
  return widths.map(w=>{const cx=x+w/2;x+=w+g;return cx;});
}

function treeDepth(nodes,cur=0){if(!nodes||!nodes.length)return cur;return Math.max(...nodes.map(n=>treeDepth(n.children||[],cur+1)));}
function getLevelH(nodes,level,cur=0){let m=BH_MIN;nodes.forEach(n=>{if(cur===level)m=Math.max(m,boxH(n.name));if(n.children&&cur<level)m=Math.max(m,getLevelH(n.children,level,cur+1));});return m;}
function getLevelW(nodes,level,cur=0){let m=BW_MIN;nodes.forEach(n=>{if(cur===level)m=Math.max(m,boxW(n.name));if(n.children&&cur<level)m=Math.max(m,getLevelW(n.children,level,cur+1));});return m;}
function normDims(nodes,lw,lh,cur=0){nodes.forEach(n=>{n._w=lw[cur];n._h=lh[cur];if(n.children)normDims(n.children,lw,lh,cur+1);});}
function assignCompact(nodes,center,gap){
  const ws=nodes.map(n=>n._w);
  const cxs=calcPos(ws,gap,center);
  nodes.forEach((n,i)=>{n._cx=cxs[i];if(n.children&&n.children.length)assignCompact(n.children,cxs[i],gap);});
}

function renderPreview(){
  // Read spacing fresh every time
  const SH_GAP=spacing.shGap, PC_GAP_=spacing.pcGap, SUB_GAP=spacing.subGap, GAP_V=spacing.gapV;

  const comp=document.getElementById('companyName').value||'本公司';
  const compFilled=document.getElementById('companyFilled').checked;
  const shs=shareholders, subs=JSON.parse(JSON.stringify(subsidiaries)), mid=middleTier;
  const n=shs.length;

  // Shareholder sizes
  const sh_w=shs.map(s=>boxW(s.name)), sh_h=shs.map(s=>boxH(s.name));
  const msw=Math.max(...sh_w,BW_MIN), msh=Math.max(...sh_h,BH_MIN);
  const sh_cx=calcPos(Array(n).fill(msw),SH_GAP);
  // Snap to center
  for(let i=0;i<sh_cx.length;i++) if(Math.abs(sh_cx[i]-SW/2)<0.6) sh_cx[i]=SW/2;

  const comp_w=boxW(comp), comp_h=boxH(comp);
  const PC_H_=0.85, PC_STUB_=0.4, PC_BUS_=0.5;
  const pcAbove=(parentControllers&&parentControllers.length)?(PC_H_+PC_STUB_+PC_BUS_):0;

  // Estimate total height for vertical centering
  const dp_pre=treeDepth(subs);
  const lhPre=Array.from({length:Math.max(dp_pre,1)},(_,lv)=>getLevelH(subs,lv));
  let hCont=pcAbove+msh+GAP_V;
  if(mid.length){hCont+=GAP_V+Math.max(...mid.map(m=>boxH(m.name)),BH_MIN)+GAP_V;}
  hCont+=comp_h+GAP_V;
  if(subs.length){hCont+=GAP_V+lhPre.reduce((a,lh)=>a+lh+GAP_V*2,0);}
  const R1Y=Math.max(pcAbove+0.8,(SH-hCont)/2+pcAbove);
  const CCX=SW/2;

  const wrap=document.getElementById('svgWrap');
  const sc=40;
  const p=v=>+(v*sc).toFixed(1);

  const el=[];
  el.push('<defs><marker id="tri" markerWidth="7" markerHeight="7" refX="7" refY="3.5" orient="auto"><polygon points="0 0,7 3.5,0 7" fill="#000"/></marker></defs>');

  let _ci=0;
  const rect=(x,y,w,h,sw=1,dash='',fill='none')=>
    `<rect x="${p(x)}" y="${p(y)}" width="${p(w)}" height="${p(h)}" fill="${fill}" stroke="#000" stroke-width="${sw}" ${dash?`stroke-dasharray="${dash}"`:''}/>`;
  const ln=(x1,y1,x2,y2,sw=0.7)=>
    `<line x1="${p(x1)}" y1="${p(y1)}" x2="${p(x2)}" y2="${p(y2)}" stroke="#000" stroke-width="${sw}" shape-rendering="crispEdges"/>`;
  const arr=(x,y1,y2)=>
    `<line x1="${p(x)}" y1="${p(y1)}" x2="${p(x)}" y2="${p(y2)}" stroke="#000" stroke-width="0.7" marker-end="url(#tri)" shape-rendering="crispEdges"/>`;
  const harr=(x1,y,x2)=>
    `<line x1="${p(x1)}" y1="${p(y)}" x2="${p(x2)}" y2="${p(y)}" stroke="#000" stroke-width="0.7" marker-end="url(#tri)" shape-rendering="crispEdges"/>`;
  const txt=(x,y,w,h,lines,fs=8)=>{
    const id='cl'+(++_ci);
    const lh=fs*1.5,tot=lines.length*lh,sy=p(y)+p(h)/2-tot/2+fs;
    return `<clipPath id="${id}"><rect x="${p(x+0.05)}" y="${p(y+0.05)}" width="${p(w-0.1)}" height="${p(h-0.1)}"/></clipPath>`
      +lines.map((l,i)=>`<text x="${p(x)+p(w)/2}" y="${+(sy+i*lh).toFixed(1)}" text-anchor="middle" font-size="${fs}" font-family="serif" clip-path="url(#${id})" fill="#000">${esc(l)}</text>`).join('');
  };
  const lbl=(x,y,text)=>
    `<text x="${p(x)}" y="${p(y)+p(PH)/2+3}" text-anchor="start" font-size="7" font-family="serif" fill="#555">${esc(text)}</text>`;

  // Dashed control box
  const ctrl=shs.map((s,i)=>s.isControl?i:-1).filter(i=>i>=0);
  if(ctrl.length){
    const lx=sh_cx[ctrl[0]]-msw/2-0.3, rx=sh_cx[ctrl[ctrl.length-1]]+msw/2+0.3;
    el.push(rect(lx,R1Y-0.3,rx-lx,msh+0.6,0.5,'4,3'));
  }

  // Parent controllers (T-connector)
  const BUS_Y=R1Y-PC_BUS_;
  const BOX_BOT=BUS_Y-PC_STUB_;
  if(parentControllers&&parentControllers.length){
    parentControllers.forEach(pc=>{
      const pc_w=Math.max(BW_MIN,boxW(pc.name));
      const ctrlIdxs=(pc.shareholders||[]).map(s=>+s.idx).filter(idx=>idx<sh_cx.length);
      if(!ctrlIdxs.length)return;
      const cx=(Math.min(...ctrlIdxs.map(i=>sh_cx[i]))+Math.max(...ctrlIdxs.map(i=>sh_cx[i])))/2;
      el.push(rect(cx-pc_w/2,BOX_BOT-PC_H_,pc_w,PC_H_,1.2));
      el.push(txt(cx-pc_w/2,BOX_BOT-PC_H_,pc_w,PC_H_,pc.name.split(/\\n/)));
      el.push(ln(cx,BOX_BOT,cx,BUS_Y));
      const xs=ctrlIdxs.map(i=>sh_cx[i]);
      if(xs.length>1) el.push(ln(Math.min(...xs),BUS_Y,Math.max(...xs),BUS_Y));
      ctrlIdxs.forEach((idx,j)=>{
        const pct=(pc.shareholders[j]||{}).pct||'';
        el.push(arr(sh_cx[idx],BUS_Y,R1Y));
        if(pct) el.push(lbl(sh_cx[idx]+POFF,BUS_Y+(PC_BUS_-PH)/2,pct));
      });
    });
  }

  // Shareholders
  shs.forEach((sh,i)=>{
    el.push(rect(sh_cx[i]-msw/2,R1Y,msw,msh,1.2));
    el.push(txt(sh_cx[i]-msw/2,R1Y,msw,msh,sh.name.split(/\\n/)));
  });

  let cur_y=R1Y+msh, CY;

  // Middle tier
  if(mid.length){
    const mid_w=mid.map(m=>boxW(m.name)), mid_h=mid.map(m=>boxH(m.name));
    const mmid=Math.max(...mid_h,BH_MIN), mmw=Math.max(...mid_w,BW_MIN);
    const mid_cx=calcPos(mid_w,SH_GAP);
    const CONN_Y=cur_y+GAP_V*0.5, MID_Y=CONN_Y+GAP_V*0.5;
    const MID_BOT=MID_Y+mmid;
    CY=MID_BOT+GAP_V;
    const midPar=new Set(mid.flatMap(m=>m.parentShareholders||[]));
    mid.forEach((m,i)=>{
      el.push(rect(mid_cx[i]-mmw/2,MID_Y,mmw,mmid,1.2));
      el.push(txt(mid_cx[i]-mmw/2,MID_Y,mmw,mmid,m.name.split(/\\n/)));
      // mid-tier bottom → CY merge
      el.push(ln(mid_cx[i],MID_BOT,mid_cx[i],CY));
      el.push(lbl(mid_cx[i]+POFF,MID_BOT+(CY-MID_BOT)/2-PH/2,m.mainPct||''));
      // Parent shareholders T-junction
      const pis=(m.parentShareholders||[]).filter(pi=>pi<sh_cx.length);
      if(pis.length){
        const xs=pis.map(pi=>sh_cx[pi]);
        pis.forEach(pi=>{
          el.push(ln(sh_cx[pi],cur_y,sh_cx[pi],CONN_Y));
          el.push(lbl(sh_cx[pi]+POFF,cur_y+(CONN_Y-cur_y)/2-PH/2,shs[pi].pct||''));
        });
        if(xs.length>1) el.push(ln(Math.min(...xs),CONN_Y,Math.max(...xs),CONN_Y));
        const busX=(Math.min(...xs)+Math.max(...xs))/2;
        if(Math.abs(busX-mid_cx[i])>0.05) el.push(ln(busX,CONN_Y,mid_cx[i],CONN_Y));
        el.push(arr(mid_cx[i],CONN_Y,MID_Y));
      }
    });
    // Non-participating shareholders: direct line to CY merge
    shs.forEach((sh,i)=>{
      if(!midPar.has(i)){
        el.push(ln(sh_cx[i],cur_y,sh_cx[i],CY));
        el.push(lbl(sh_cx[i]+POFF,cur_y+(CY-cur_y)/2-PH/2,sh.pct||''));
      }
    });
    // Horizontal merge bus at CY
    const allX=[...shs.map((_,i)=>midPar.has(i)?null:sh_cx[i]),...mid_cx].filter(x=>x!==null);
    if(allX.length>1) el.push(ln(Math.min(...allX),CY,Math.max(...allX),CY));
    el.push(arr(CCX,CY,CY+GAP_V));
    CY=CY+GAP_V;
  } else {
    const MERGE_Y=cur_y+GAP_V; CY=MERGE_Y+GAP_V;
    shs.forEach((sh,i)=>{el.push(ln(sh_cx[i],cur_y,sh_cx[i],MERGE_Y));el.push(lbl(sh_cx[i]+POFF,cur_y+(MERGE_Y-cur_y)/2-PH/2,sh.pct));});
    if(n>1) el.push(ln(sh_cx[0],MERGE_Y,sh_cx[n-1],MERGE_Y));
    el.push(arr(CCX,MERGE_Y,CY));
  }

  // Company box
  const cfill=compFilled?'#D9D9D9':'none';
  el.push(rect(CCX-comp_w/2,CY,comp_w,comp_h,1.2,'',cfill));
  el.push(txt(CCX-comp_w/2,CY,comp_w,comp_h,[comp]));

  if(!subs.length){
    const svgW=Math.max(SW*sc,wrap.clientWidth-4);
    const _nsAllX=[...sh_cx.map(x=>x-msw/2),...sh_cx.map(x=>x+msw/2),CCX-comp_w/2,CCX+comp_w/2];
    const _nsL=Math.min(..._nsAllX)-0.6, _nsR=Math.max(..._nsAllX)+0.6;
    const vbH_n=+(CY+comp_h+1.5)*sc;
    const vbL_n=_nsL*sc, vbW_n=(_nsR-_nsL)*sc;
    const fitScale_n=Math.min(1,(wrap.clientWidth-8)/vbW_n,(wrap.clientHeight||600)/vbH_n);
    const zoom_n=previewZoom>0?previewZoom:fitScale_n;
    const dW_n=Math.round(vbW_n*zoom_n), dH_n=Math.round(vbH_n*zoom_n);
    wrap.innerHTML=`<svg viewBox="${vbL_n.toFixed(1)} 0 ${vbW_n.toFixed(1)} ${vbH_n.toFixed(0)}" width="${dW_n}" height="${dH_n}" xmlns="http://www.w3.org/2000/svg">${el.join('')}</svg>`;
    return;
  }

  // Subsidiaries
  const depth=treeDepth(subs);
  const levelH=Array.from({length:Math.max(depth,1)},(_,lv)=>getLevelH(subs,lv));
  const levelW=Array.from({length:Math.max(depth,1)},(_,lv)=>getLevelW(subs,lv));
  normDims(subs,levelW,levelH);
  assignCompact(subs,SW/2,SUB_GAP);

  const SM=CY+comp_h+GAP_V;
  const levelY=[SM+GAP_V];
  for(let lv=1;lv<depth;lv++) levelY.push(levelY[lv-1]+levelH[lv-1]+GAP_V*2);

  el.push(ln(CCX,CY+comp_h,CCX,SM));
  if(subs.length>1){
    const lx=Math.min(...subs.map(s=>s._cx)), rx=Math.max(...subs.map(s=>s._cx));
    el.push(ln(lx,SM,rx,SM));
  }

  const boxPos={};
  function drawTree(nodes,parentSM,level){
    const mh=levelH[level], mw=levelW[level], SY=levelY[level];
    nodes.forEach(node=>{
      const cx=node._cx;
      el.push(arr(cx,parentSM,SY));
      el.push(lbl(cx+POFF,parentSM+(SY-parentSM)/2-PH/2,node.pct));
      el.push(rect(cx-mw/2,SY,mw,mh,1.2));
      el.push(txt(cx-mw/2,SY,mw,mh,node.name.split(/\\n/)));
      boxPos[node.name.replace(/\\\\n/g,'')]={cx,w:mw,y:SY,h:mh};
      if(node.children&&node.children.length){
        const ch=node.children;
        const chSM=SY+mh+GAP_V;
        el.push(ln(cx,SY+mh,cx,chSM));
        if(ch.length>1) el.push(ln(Math.min(...ch.map(c=>c._cx)),chSM,Math.max(...ch.map(c=>c._cx)),chSM));
        drawTree(ch,chSM,level+1);
      }
    });
  }
  drawTree(subs,SM,0);

  // Cross-links (horizontal arrows)
  const nl_=new RegExp(String.fromCharCode(92,110),'g');
  crossLinks.forEach(xl=>{
    const fp=boxPos[(xl.from||'').replace(nl_,'')];
    const tp=boxPos[(xl.to||'').replace(nl_,'')];
    if(fp&&tp){
      const midY=(fp.y+fp.h/2+tp.y+tp.h/2)/2;
      const x1=fp.cx<tp.cx?fp.cx+fp.w/2:fp.cx-fp.w/2;
      const x2=fp.cx<tp.cx?tp.cx-tp.w/2:tp.cx+tp.w/2;
      el.push(harr(x1,midY,x2));
      if(xl.pct) el.push(`<text x="${p((x1+x2)/2)}" y="${p(midY)-12}" text-anchor="middle" font-size="7" font-family="serif" fill="#555">${esc(xl.pct)}</text>`);
    }
  });

  // Compute SVG dimensions
  let maxH=SM+GAP_V;
  levelY.forEach((ly,lv)=>{maxH=Math.max(maxH,ly+(levelH[lv]||BH_MIN)+0.8);});
  const allRX=[...subs.map(s=>s._cx+s._w/2),CCX+comp_w/2,...sh_cx.map(x=>x+msw/2)];
  const allLX=[...subs.map(s=>s._cx-s._w/2),CCX-comp_w/2,...sh_cx.map(x=>x-msw/2)];
  const contentL=Math.min(...allLX)-0.6, contentR=Math.max(...allRX)+0.6;
  const contentW=contentR-contentL;
  const vbL_s=contentL*sc, vbW_s=contentW*sc, vbH_s=maxH*sc;
    const fitScale_s=Math.min(1,(wrap.clientWidth-8)/vbW_s,(wrap.clientHeight||600)/vbH_s);
    const zoom_s=previewZoom>0?previewZoom:fitScale_s;
    const dW_s=Math.round(vbW_s*zoom_s), dH_s=Math.round(vbH_s*zoom_s);
    wrap.innerHTML=`<svg viewBox="${vbL_s.toFixed(1)} 0 ${vbW_s.toFixed(1)} ${vbH_s.toFixed(1)}" width="${dW_s}" height="${dH_s}" xmlns="http://www.w3.org/2000/svg">${el.join('')}</svg>`;
}

// ── Generate PPTX ─────────────────────────────────────────────────────────────
async function generate(){
  const btn=document.getElementById('genBtn'),st=document.getElementById('status');
  btn.disabled=true; st.textContent='正在生成...'; st.className='';
  try{
    const body={
      companyName:document.getElementById('companyName').value||'本公司',
      companyFilled:document.getElementById('companyFilled').checked,
      shareholders,
      parentControllers,
      middleTier,
      subsidiaries,
      crossLinks,
      divider:null,
      spacing,
    };
    const resp=await fetch('/generate',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(body)});
    if(!resp.ok) throw new Error(await resp.text());
    const blob=await resp.blob();
    const url=URL.createObjectURL(blob);
    const a=document.createElement('a'); a.href=url;
    const enc=encodeURIComponent('股权架构图.pptx');
    a.download='股权架构图.pptx';
    a.setAttribute('download','股权架构图.pptx');
    document.body.appendChild(a); a.click(); document.body.removeChild(a);
    URL.revokeObjectURL(url);
    st.textContent='文件已生成并下载！'; st.className='ok';
  }catch(e){st.textContent='生成失败：'+e.message; st.className='err';}
  btn.disabled=false;
}

// ── Resizable panel ─────────────────────────────────────────────────────────────
(function(){
  const handle=document.getElementById('resizeHandle');
  const panel=document.getElementById('mainPanel');
  if(!handle||!panel)return;
  let dragging=false,startX=0,startW=0;
  handle.addEventListener('mousedown',e=>{
    dragging=true; startX=e.clientX; startW=panel.offsetWidth;
    handle.classList.add('dragging');
    document.body.style.cursor='col-resize';
    document.body.style.userSelect='none';
  });
  document.addEventListener('mousemove',e=>{
    if(!dragging)return;
    const w=Math.max(260,Math.min(700,startW+(e.clientX-startX)));
    panel.style.width=w+'px';
  });
  document.addEventListener('mouseup',()=>{
    if(!dragging)return;
    dragging=false; handle.classList.remove('dragging');
    document.body.style.cursor=''; document.body.style.userSelect='';
    renderPreview();
  });
})();

// ── Init ──────────────────────────────────────────────────────────────────────
renderShareholders();
renderParentControllers();
renderMid();
renderSubsidiariesUI(subsidiaries,null,document.getElementById('subsidiaryList'));
renderXL();
renderSpacing();
// 直接调用（内联脚本无需等待load事件）
function safeRenderPreview(){
  try { renderPreview(); }
  catch(e){
    var w=document.getElementById('svgWrap');
    if(w) w.textContent='预览出错: '+e.message;
    console.error('renderPreview error:',e);
  }
}
setTimeout(safeRenderPreview, 0);
window.addEventListener('resize', safeRenderPreview);
</script>
</body>
</html>"""

class Handler(BaseHTTPRequestHandler):
    def log_message(self, fmt, *args): pass
    def do_GET(self):
        self.send_response(200)
        self.send_header('Content-Type','text/html; charset=utf-8')
        self.end_headers()
        self.wfile.write(HTML.encode('utf-8'))
    def do_POST(self):
        if self.path != '/generate':
            self.send_response(404); self.end_headers(); return
        try:
            length = int(self.headers.get('Content-Length',0))
            data = json.loads(self.rfile.read(length))
            pptx = generate_pptx(data)
            name_enc = "UTF-8''%E8%82%A1%E6%9D%83%E6%9E%B6%E6%9E%84%E5%9B%BE.pptx"
            self.send_response(200)
            self.send_header('Content-Type','application/vnd.openxmlformats-officedocument.presentationml.presentation')
            self.send_header('Content-Disposition', f"attachment; filename*={name_enc}")
            self.send_header('Content-Length', str(len(pptx)))
            self.end_headers()
            self.wfile.write(pptx)
        except Exception as e:
            import traceback
            msg = traceback.format_exc().encode()
            self.send_response(500)
            self.send_header('Content-Type','text/plain')
            self.send_header('Content-Length',str(len(msg)))
            self.end_headers()
            self.wfile.write(msg)

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5001))
    print(f"Starting on http://0.0.0.0:{port}")
    HTTPServer(('0.0.0.0', port), Handler).serve_forever()
